# -*- coding: utf-8 -*-
"""
ROI Algorithm Simulation Results PPT Generator
Creates a presentation summarizing the ROI processing pipeline results
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
from PIL import Image

# Paths
work_dir = Path(r'c:\Users\byungpaul\Desktop\AI_Project\20260304_ROI_algorithm')
output_dir = work_dir / 'output'
ppt_path = work_dir / 'ROI_Algorithm_Simulation_Results_EN.pptx'

# Create presentation (16:9 aspect ratio)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Define layouts
title_slide_layout = prs.slide_layouts[6]  # Blank
content_slide_layout = prs.slide_layouts[6]  # Blank


def get_image_size_to_fit(image_path, max_width_inches, max_height_inches):
    """Calculate image size to fit within bounds while maintaining aspect ratio"""
    try:
        with Image.open(image_path) as img:
            img_width, img_height = img.size
            aspect_ratio = img_width / img_height

            # Calculate dimensions that fit within bounds
            if aspect_ratio > (max_width_inches / max_height_inches):
                # Image is wider - constrain by width
                width = max_width_inches
                height = width / aspect_ratio
            else:
                # Image is taller - constrain by height
                height = max_height_inches
                width = height * aspect_ratio

            return Inches(width), Inches(height)
    except Exception:
        return Inches(max_width_inches), None


def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide = prs.slides.add_slide(title_slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, content_items=None, image_path=None, two_images=None):
    """Add a content slide with optional bullet points and/or image"""
    slide = prs.slides.add_slide(content_slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True

    # Add horizontal line under title
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(12.333), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0, 120, 215)
    line.line.fill.background()

    # Available space for content
    content_top = 1.3
    content_height = 5.8  # Leave room for bottom margin

    if two_images:
        # Two images side by side
        img1_path, img2_path = two_images
        max_img_width = 5.8
        max_img_height = content_height

        if Path(img1_path).exists():
            w, h = get_image_size_to_fit(img1_path, max_img_width, max_img_height)
            slide.shapes.add_picture(str(img1_path), Inches(0.5), Inches(content_top), width=w, height=h)
        if Path(img2_path).exists():
            w, h = get_image_size_to_fit(img2_path, max_img_width, max_img_height)
            slide.shapes.add_picture(str(img2_path), Inches(7), Inches(content_top), width=w, height=h)
    elif image_path and content_items:
        # Image on right, content on left
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(content_top), Inches(5.5), Inches(content_height))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(16)
            p.space_after = Pt(8)

        if Path(image_path).exists():
            max_img_width = 6.5
            max_img_height = content_height
            w, h = get_image_size_to_fit(image_path, max_img_width, max_img_height)
            slide.shapes.add_picture(str(image_path), Inches(6.3), Inches(content_top), width=w, height=h)
    elif image_path:
        # Full width image - fit to available space
        if Path(image_path).exists():
            max_img_width = 12.333
            max_img_height = content_height
            w, h = get_image_size_to_fit(image_path, max_img_width, max_img_height)
            # Center the image horizontally
            left = (13.333 - w.inches) / 2
            slide.shapes.add_picture(str(image_path), Inches(left), Inches(content_top), width=w, height=h)
    elif content_items:
        # Content only
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(content_top), Inches(12.333), Inches(content_height))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(20)
            p.space_after = Pt(14)

    return slide


def add_image_slide(prs, title, image_path, caption=""):
    """Add a slide with centered image that fits within bounds"""
    slide = prs.slides.add_slide(content_slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True

    # Image - fit within available space
    content_top = 1.2
    content_height = 5.2 if caption else 5.8
    max_img_width = 12.333

    if Path(image_path).exists():
        w, h = get_image_size_to_fit(image_path, max_img_width, content_height)
        # Center the image horizontally
        left = (13.333 - w.inches) / 2
        slide.shapes.add_picture(str(image_path), Inches(left), Inches(content_top), width=w, height=h)

    # Caption
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(12.333), Inches(0.5))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(14)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER

    return slide


# ============== Create Slides ==============

# Slide 1: Title
add_title_slide(
    prs,
    "Display Panel ROI Algorithm",
    "Simulation Results & Analysis\n\nMR Display Hardware Team | March 2026"
)

# Slide 2: Overview
add_content_slide(
    prs,
    "Project Overview",
    content_items=[
        "Objective: Extract display pixel values from camera-captured images",
        "Challenge: Camera resolution >> Display resolution (approx. 3.9:1 ratio)",
        "Solution: Area-sum based pixel extraction with ROI detection",
        "Input: 14144 x 10640 pixel camera image (16-bit)",
        "Output: 2412 x 2288 pixel display image (16-bit normalized)",
        "Key Features:",
        "    - Automatic ROI (Region of Interest) detection",
        "    - Perspective warp correction for tilt compensation",
        "    - Area-sum resizing to preserve total light intensity",
        "    - 16-bit normalization for maximum dynamic range"
    ]
)

# Slide 3: Pipeline Overview
add_content_slide(
    prs,
    "Processing Pipeline",
    content_items=[
        "Step 1: Image Loading - Load 16-bit TIFF camera image",
        "Step 2: Image Cropping - Remove unnecessary border regions",
        "Step 3: ROI Detection - Detect display panel boundaries",
        "Step 4: Corner Extraction - Find 4 corner points precisely",
        "Step 5: Perspective Warp - Correct tilt/rotation",
        "Step 6: Area-Sum Resize - Map camera pixels to display pixels",
        "Step 7: 16-bit Normalization - Scale to full dynamic range"
    ],
    image_path=output_dir / "06_final_comparison.png"
)

# Slide 4: Step 1-2 - Image Loading & Cropping
add_image_slide(
    prs,
    "Step 1-2: Image Loading & Cropping",
    output_dir / "00_crop_comparison.png",
    "Original: 14144x10640 → Cropped: 10000x9600 pixels"
)

# Slide 5: Step 3 - ROI Detection
add_content_slide(
    prs,
    "Step 3: ROI Detection",
    content_items=[
        "Algorithm: Threshold-based segmentation + morphological operations",
        "Detected ROI: 9374 x 8894 pixels",
        "Tilt angle: -0.03 degrees",
        "Contour area: 83,316,360 pixels",
        "Method:",
        "    - Binary thresholding (threshold=50)",
        "    - Morphological closing (kernel=51x51)",
        "    - Contour detection and filtering",
        "    - MinAreaRect for rotation estimation"
    ],
    image_path=output_dir / "02_roi_detection.png"
)

# Slide 6: Step 4 - Corner Detection
add_image_slide(
    prs,
    "Step 4: Corner Detection",
    output_dir / "03_corner_detection.png",
    "Precise 4-corner extraction using edge crossing method"
)

# Slide 7: Corner Detection Detail
add_image_slide(
    prs,
    "Corner Detection - Edge Analysis",
    output_dir / "07_edge_corners_star_marks.png",
    "Edge profile analysis with star markers at detected corners"
)

# Slide 8: Step 5 - Perspective Warp
add_image_slide(
    prs,
    "Step 5: Perspective Warp Correction",
    output_dir / "04_warp_comparison.png",
    "INTER_NEAREST interpolation to preserve original pixel values"
)

# Slide 9: Step 6 - Area Sum Resize
add_content_slide(
    prs,
    "Step 6: Area-Sum Resize",
    content_items=[
        "Purpose: Map multiple camera pixels to single display pixel",
        "Scale factors: 3.8864 x 3.8872 (camera/display ratio)",
        "Input: 9374 x 8894 pixels (warped)",
        "Output: 2412 x 2288 pixels (display resolution)",
        "Algorithm: Weighted area summation",
        "    - Each display pixel = sum of corresponding camera pixels",
        "    - Fractional boundaries handled with proportional weights",
        "    - Preserves total light intensity",
        "Processing time: ~67 seconds"
    ],
    image_path=output_dir / "05_resize_comparison.png"
)

# Slide 10: Pixel Mapping Visualization
add_image_slide(
    prs,
    "Area-Sum Resize - Pixel Mapping",
    output_dir / "19_4x4_box_mapping.png",
    "Camera pixels (red boxes) → Display pixels (output grid)"
)

# Slide 11: 3x3 Block Analysis
add_image_slide(
    prs,
    "3x3 Block Analysis",
    output_dir / "3x3_max_extraction_result.png",
    "3x3 block max extraction for display pixel localization"
)

# Slide 12: Step 7 - Normalization
add_content_slide(
    prs,
    "Step 7: 16-bit Normalization",
    content_items=[
        "Input range: 283.67 ~ 14,511.35",
        "Output range: 0 ~ 65,535 (full 16-bit)",
        "Purpose: Maximize dynamic range utilization",
        "Formula: normalized = (value - min) / (max - min) × 65535",
        "Output format: 16-bit TIFF"
    ],
    image_path=output_dir / "05_final_result.png"
)

# Slide 13: Quality Validation
add_image_slide(
    prs,
    "Quality Validation",
    output_dir / "quality_validation_charts.png",
    "Statistical analysis of processing quality"
)

# Slide 14: Detailed Comparison
add_image_slide(
    prs,
    "Input vs Output Comparison",
    output_dir / "16_input_output_comparison.png",
    "Side-by-side comparison of processing stages"
)

# Slide 15: Pixel Level Detail
add_image_slide(
    prs,
    "Pixel Level Analysis",
    output_dir / "18_pixel_level_detail.png",
    "Detailed pixel-level verification"
)

# Slide 16: Results Summary
add_content_slide(
    prs,
    "Results Summary",
    content_items=[
        "✓ Successfully extracted display panel ROI from camera image",
        "✓ Achieved accurate corner detection (sub-pixel precision)",
        "✓ Perspective correction with minimal distortion",
        "✓ Area-sum resize preserves total intensity",
        "✓ Full 16-bit dynamic range utilization",
        "",
        "Performance Metrics:",
        "    - Input: 14144 × 10640 pixels (16-bit)",
        "    - Output: 2412 × 2288 pixels (16-bit normalized)",
        "    - Scale ratio: 3.89:1 (matches display pitch)",
        "    - Processing time: < 70 seconds",
        "",
        "Output Files:",
        "    - 06_final_normalized.tif (main output)",
        "    - 06_final_normalized_preview.png (8-bit preview)"
    ]
)

# Slide 17: Code Architecture
add_content_slide(
    prs,
    "Code Architecture",
    content_items=[
        "Modular Python implementation with 7 core modules:",
        "",
"1_utils.py - Utility functions (load, save, normalize)",
        "2_roi_detector.py - ROI detection with adaptive thresholding",
        "3_image_cropper.py - Image cropping with region definition",
        "4_perspective_warper.py - Perspective transformation",
        "5_area_sum_resizer.py - Area-sum based resizing",
        "6_image_normalizer.py - Bit-depth normalization",
        "7_display_panel_processor.py - Main pipeline orchestrator",
        "",
        "Features:",
        "    - Configurable via ProcessingConfig dataclass",
        "    - Progress tracking with verbose output",
        "    - Intermediate result saving for debugging"
    ]
)

# Slide 18: Next Steps
add_content_slide(
    prs,
    "Next Steps & Future Work",
    content_items=[
        "Performance Optimization:",
        "    - Numba JIT compilation for resize operations",
        "    - Parallel processing for multi-image batches",
        "",
        "Algorithm Improvements:",
        "    - Sub-pixel corner detection refinement",
        "    - Adaptive threshold optimization",
        "    - Multi-frame averaging for noise reduction",
        "",
        "Validation:",
        "    - Cross-validation with reference measurements",
        "    - Statistical analysis of pixel accuracy",
        "",
        "Integration:",
        "    - Automated calibration pipeline",
        "    - Real-time processing capability"
    ]
)

# Slide 19: Thank You
add_title_slide(
    prs,
    "Thank You",
    "Questions & Discussion\n\nContact: byungpaul@meta.com"
)

# Save presentation
prs.save(str(ppt_path))
print(f"Presentation saved to: {ppt_path}")
print(f"Total slides: {len(prs.slides)}")
