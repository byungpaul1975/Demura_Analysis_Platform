# -*- coding: utf-8 -*-
"""
ROI Pipeline PPT Generator — Step-by-step with ALL intermediate results.

Runs the full processing pipeline in code order:
  1_utils → 3_image_cropper → 2_roi_detector → 4_perspective_warper
  → 5_area_sum_resizer → 6_image_normalizer

Generates intermediate visualisation images for every step and builds
a 16:9 English PowerPoint presentation.

Usage:
    python docs/create_pipeline_ppt.py
"""
import sys
import numpy as np
import cv2
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pathlib import Path

PROJECT_ROOT = Path(__file__).parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

from src import (
    ImageCropper, CropRegion,
    ROIDetector, AdaptiveROIDetector,
    PerspectiveWarper,
    AreaSumResizer,
    ImageNormalizer,
    ProcessingConfig,
)

# ── paths ──────────────────────────────────────────────────────────
DOCS_DIR  = Path(__file__).parent
IMG_DIR   = DOCS_DIR / "pipeline_images"
IMG_DIR.mkdir(parents=True, exist_ok=True)

DATA_DIR  = PROJECT_ROOT / "data"
IMAGE_FILE = DATA_DIR / "G32_cal.tif"

# ── colours ────────────────────────────────────────────────────────
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BLUE    = RGBColor(0x00, 0x3C, 0x71)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0x33, 0x33, 0x33)
LTBLUE  = RGBColor(0xBB, 0xDE, 0xFB)
MIDBLUE = RGBColor(0x90, 0xCA, 0xF9)
DKBLUE  = RGBColor(0x64, 0xB5, 0xF6)


# ═══════════════════════════════════════════════════════════════════
# Helper: save matplotlib figure
# ═══════════════════════════════════════════════════════════════════
def _save(fig, name, dpi=180):
    path = IMG_DIR / name
    fig.savefig(str(path), dpi=dpi, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    print(f"  saved -> {path.name}")
    return str(path)


def _to8(img):
    """Convert any image to uint8 for visualisation."""
    f = img.astype(np.float64)
    mn, mx = f.min(), max(f.max(), 1)
    return np.clip((f - mn) / (mx - mn) * 255, 0, 255).astype(np.uint8)


# ═══════════════════════════════════════════════════════════════════
# Pipeline execution + image generation
# ═══════════════════════════════════════════════════════════════════
def run_pipeline():
    """Run the full pipeline and generate all intermediate images."""
    config = ProcessingConfig(
        crop_x_start=1700, crop_y_start=300,
        crop_x_end=11700, crop_y_end=9900,
        use_crop=True,
        roi_threshold=50, morph_kernel_size=51,
        display_width=2412, display_height=2288,
        output_bit_depth=16, save_intermediates=True,
    )
    images = {}   # name → filepath

    # ── Step 0  Load original image (1_utils.py) ──────────────────
    print("\n[Step 0] Loading image …")
    original = cv2.imread(str(IMAGE_FILE), cv2.IMREAD_UNCHANGED)
    if original is None:
        raise FileNotFoundError(f"Cannot load {IMAGE_FILE}")
    print(f"  shape={original.shape}, dtype={original.dtype}, "
          f"range=[{original.min()}, {original.max()}]")

    # 0-a  original overview
    fig, ax = plt.subplots(figsize=(12, 9))
    ax.imshow(_to8(original), cmap="gray")
    ax.set_title(f"Original Image: {original.shape[1]}×{original.shape[0]}  |  "
                 f"dtype={original.dtype}  |  range=[{original.min()}, {original.max()}]",
                 fontsize=12, fontweight="bold")
    ax.set_xlabel("X (pixels)"); ax.set_ylabel("Y (pixels)")
    images["s0_original"] = _save(fig, "s0_original.png")

    # 0-b  histogram
    fig, axes = plt.subplots(1, 2, figsize=(16, 5))
    axes[0].hist(original.ravel(), bins=256, color="steelblue", alpha=0.8)
    axes[0].set_title("Pixel Value Histogram", fontweight="bold")
    axes[0].set_xlabel("Pixel Value"); axes[0].set_ylabel("Frequency")
    axes[0].axvline(x=config.roi_threshold, color="red", ls="--", lw=2,
                    label=f"Threshold={config.roi_threshold}")
    axes[0].legend()
    axes[0].set_yscale("log")
    axes[1].axis("off")
    info = (f"File      : {IMAGE_FILE.name}\n"
            f"Resolution: {original.shape[1]} × {original.shape[0]}\n"
            f"Data Type : {original.dtype}\n"
            f"Min / Max : {original.min()} / {original.max()}\n"
            f"Mean      : {original.mean():.1f}\n"
            f"Std       : {original.std():.1f}")
    axes[1].text(0.1, 0.7, info, fontsize=13, va="top", family="monospace",
                 transform=axes[1].transAxes,
                 bbox=dict(boxstyle="round", facecolor="lightyellow"))
    plt.tight_layout()
    images["s0_histogram"] = _save(fig, "s0_histogram.png")

    # ── Step 1  Image cropping (3_image_cropper.py) ───────────────
    print("\n[Step 1] Cropping image …")
    cropper = ImageCropper(CropRegion(
        config.crop_x_start, config.crop_y_start,
        config.crop_x_end, config.crop_y_end))
    cropped = cropper.crop(original)
    print(f"  cropped: {cropped.shape[1]}×{cropped.shape[0]}")

    fig, axes = plt.subplots(1, 2, figsize=(18, 7))
    axes[0].imshow(_to8(original), cmap="gray")
    rect = plt.Rectangle(
        (config.crop_x_start, config.crop_y_start),
        config.crop_x_end - config.crop_x_start,
        config.crop_y_end - config.crop_y_start,
        lw=2, ec="red", fc="none", ls="--")
    axes[0].add_patch(rect)
    axes[0].set_title(f"Original ({original.shape[1]}×{original.shape[0]})\n"
                      f"Red box = crop region", fontweight="bold")
    axes[1].imshow(_to8(cropped), cmap="gray")
    axes[1].set_title(f"Cropped ({cropped.shape[1]}×{cropped.shape[0]})\n"
                      f"Region: ({config.crop_x_start},{config.crop_y_start}) → "
                      f"({config.crop_x_end},{config.crop_y_end})", fontweight="bold")
    plt.tight_layout()
    images["s1_crop"] = _save(fig, "s1_crop.png")

    # ── Step 2  ROI detection (2_roi_detector.py) ─────────────────
    print("\n[Step 2] Detecting ROI …")
    detector = ROIDetector(threshold=config.roi_threshold,
                           morph_kernel_size=config.morph_kernel_size)
    roi = detector.detect(cropped)
    print(f"  ROI: {roi.width:.0f}×{roi.height:.0f}, "
          f"tilt={roi.angle:.4f}°, area={roi.area:,.0f}")

    # 2-a  binary mask + morphology
    gray = cropped if len(cropped.shape) == 2 else cv2.cvtColor(cropped, cv2.COLOR_BGR2GRAY)
    binary_raw = (gray > config.roi_threshold).astype(np.uint8) * 255
    mask = detector.get_binary_mask()

    fig, axes = plt.subplots(1, 3, figsize=(20, 6))
    axes[0].imshow(_to8(gray), cmap="gray")
    axes[0].set_title(f"Grayscale Input\nThreshold = {config.roi_threshold}", fontweight="bold")
    axes[1].imshow(binary_raw, cmap="gray")
    axes[1].set_title(f"Binary (pixel > {config.roi_threshold})", fontweight="bold")
    if mask is not None:
        axes[2].imshow(mask, cmap="gray")
    axes[2].set_title(f"After Morphology ({config.morph_kernel_size}×{config.morph_kernel_size})\n"
                      f"+ Hole Fill", fontweight="bold")
    plt.tight_layout()
    images["s2_binary"] = _save(fig, "s2_binary.png")

    # 2-b  contour + corners
    vis_crop = cv2.cvtColor(_to8(cropped), cv2.COLOR_GRAY2BGR) if len(_to8(cropped).shape) == 2 else _to8(cropped).copy()
    vis_contour = vis_crop.copy()
    if roi.contour is not None:
        cv2.drawContours(vis_contour, [roi.contour], -1, (0, 255, 0), 3)

    vis_corners = vis_crop.copy()
    labels = ["TL", "TR", "BR", "BL"]
    colors_bgr = [(255, 0, 0), (0, 200, 0), (0, 0, 255), (255, 165, 0)]
    colors_rgb = ["blue", "green", "red", "orange"]
    for i in range(4):
        pt = tuple(roi.corners[i].astype(int))
        cv2.circle(vis_corners, pt, 14, colors_bgr[i], -1)
        cv2.putText(vis_corners, labels[i], (pt[0]+18, pt[1]-8),
                    cv2.FONT_HERSHEY_SIMPLEX, 1.6, colors_bgr[i], 3)
        pt2 = tuple(roi.corners[(i+1) % 4].astype(int))
        cv2.line(vis_corners, pt, pt2, (0, 255, 255), 2)

    fig, axes = plt.subplots(1, 2, figsize=(18, 7))
    axes[0].imshow(cv2.cvtColor(vis_contour, cv2.COLOR_BGR2RGB))
    axes[0].set_title(f"Largest Contour (green)\nArea = {roi.area:,.0f} px²", fontweight="bold")
    axes[1].imshow(cv2.cvtColor(vis_corners, cv2.COLOR_BGR2RGB))
    axes[1].set_title(f"4 Corners + Edges\nTilt = {roi.angle:.4f}°  |  "
                      f"Size = {roi.width:.0f}×{roi.height:.0f}", fontweight="bold")
    plt.tight_layout()
    images["s2_corners"] = _save(fig, "s2_corners.png")

    # 2-c  corner coordinate table
    fig, ax = plt.subplots(figsize=(8, 5))
    ax.axis("off")
    ax.set_title("Corner Coordinates & ROI Metrics", fontweight="bold", fontsize=14, pad=15)
    header = ["Corner", "X", "Y"]
    rows = [[labels[i], f"{roi.corners[i][0]:.1f}", f"{roi.corners[i][1]:.1f}"]
            for i in range(4)]
    table = ax.table(cellText=[header]+rows, loc="center", cellLoc="center",
                     colWidths=[0.25, 0.25, 0.25])
    table.auto_set_font_size(False); table.set_fontsize(12); table.scale(1, 2)
    for i in range(5):
        for j in range(3):
            c = table[i, j]
            if i == 0:
                c.set_facecolor("#003C71"); c.set_text_props(color="white", fontweight="bold")
            else:
                c.set_facecolor("#F0F2F5" if i % 2 == 0 else "white")
    ax.text(0.5, 0.05,
            f"Width={roi.width:.1f}  Height={roi.height:.1f}  "
            f"Tilt={roi.angle:.4f}°  Area={roi.area:,.0f} px²",
            fontsize=10, transform=ax.transAxes, ha="center", family="monospace",
            bbox=dict(boxstyle="round", facecolor="lightyellow"))
    images["s2_table"] = _save(fig, "s2_table.png")

    # ── Step 3  Perspective warp (4_perspective_warper.py) ─────────
    print("\n[Step 3] Perspective warp …")
    warper = PerspectiveWarper(interpolation=cv2.INTER_NEAREST)
    warp_result = warper.warp(cropped, roi.corners)
    warped = warp_result.image
    print(f"  warped: {warp_result.width}×{warp_result.height}")

    # 3-a  before / after warp
    fig, axes = plt.subplots(1, 2, figsize=(18, 7))
    axes[0].imshow(_to8(cropped), cmap="gray")
    c = roi.corners
    for i in range(4):
        axes[0].plot([c[i][0], c[(i+1)%4][0]], [c[i][1], c[(i+1)%4][1]], "r-", lw=2)
        axes[0].plot(c[i][0], c[i][1], "ro", ms=8)
    axes[0].set_title(f"Before Warp  (tilt = {roi.angle:.2f}°)", fontweight="bold")
    axes[1].imshow(_to8(warped), cmap="gray")
    axes[1].set_title(f"After Perspective Warp\n"
                      f"{warp_result.width}×{warp_result.height}  |  INTER_NEAREST",
                      fontweight="bold")
    plt.tight_layout()
    images["s3_warp"] = _save(fig, "s3_warp.png")

    # 3-b  transform matrix details
    M = warp_result.transform_matrix
    fig, ax = plt.subplots(figsize=(10, 6))
    ax.axis("off")
    ax.set_title("Perspective Transform Details", fontweight="bold", fontsize=14, pad=15)
    txt = "Transform Matrix (3×3):\n"
    for r in range(3):
        txt += f"  [{M[r,0]:+12.6f}  {M[r,1]:+12.6f}  {M[r,2]:+12.6f}]\n"
    txt += f"\nCorner Mapping (Source → Destination):\n"
    for i, lbl in enumerate(labels):
        s = warp_result.src_corners[i]; d = warp_result.dst_corners[i]
        txt += f"  {lbl}: ({s[0]:.1f}, {s[1]:.1f}) → ({d[0]:.1f}, {d[1]:.1f})\n"
    txt += f"\nInterpolation: INTER_NEAREST (preserves original pixel values)\n"
    txt += f"Output size:   {warp_result.width} × {warp_result.height}"
    ax.text(0.5, 0.5, txt, fontsize=11, ha="center", va="center", family="monospace",
            transform=ax.transAxes,
            bbox=dict(boxstyle="round,pad=0.8", facecolor="#E3F2FD"))
    images["s3_matrix"] = _save(fig, "s3_matrix.png")

    # ── Step 4  Area-sum resize (5_area_sum_resizer.py) ───────────
    print("\n[Step 4] Area-sum resize …")
    resizer = AreaSumResizer(show_progress=True)
    target = (config.display_width, config.display_height)
    resize_result = resizer.resize(warped, target)
    resized = resize_result.image
    print(f"  resized: {config.display_width}×{config.display_height}, "
          f"scale={resize_result.scale_x:.4f}×{resize_result.scale_y:.4f}")

    # 4-a  before / after resize
    fig, axes = plt.subplots(1, 2, figsize=(18, 7))
    axes[0].imshow(_to8(warped), cmap="gray")
    axes[0].set_title(f"Warped (Camera Resolution)\n"
                      f"{warp_result.width}×{warp_result.height}", fontweight="bold")
    axes[1].imshow(_to8(resized), cmap="gray")
    axes[1].set_title(f"Resized (Display Resolution)\n"
                      f"{config.display_width}×{config.display_height}", fontweight="bold")
    plt.tight_layout()
    images["s4_resize"] = _save(fig, "s4_resize.png")

    # 4-b  concept diagram
    fig, ax = plt.subplots(figsize=(12, 6))
    ax.set_xlim(0, 16); ax.set_ylim(0, 10); ax.set_aspect("equal"); ax.axis("off")
    ax.set_title("Area-Sum Resize Concept", fontweight="bold", fontsize=14)
    for i in range(6):
        ax.plot([1+i*1.5, 1+i*1.5], [1, 7.5], "k-", lw=0.5, alpha=0.5)
    for j in range(5):
        ax.plot([1, 8.5], [1+j*1.6, 1+j*1.6], "k-", lw=0.5, alpha=0.5)
    ax.add_patch(plt.Rectangle((1, 2.6), 7.5, 3.2, lw=3, ec="red", fc="red", alpha=0.12))
    ax.text(4.8, 0.3, "Camera Pixel Grid", fontsize=12, ha="center", fontweight="bold")
    ax.text(4.8, 8.8, f"1 Display Pixel ≈ {resize_result.scale_x:.2f} × {resize_result.scale_y:.2f} camera px",
            fontsize=11, ha="center", color="red", fontweight="bold")
    ax.annotate("MEAN", xy=(10, 4.2), xytext=(11.5, 4.2), fontsize=13, fontweight="bold",
                color="blue", arrowprops=dict(arrowstyle="->", color="blue", lw=2))
    ax.add_patch(plt.Rectangle((12, 3.2), 2, 2, lw=2, ec="blue", fc="blue", alpha=0.15))
    ax.text(13, 4.2, "1 px", fontsize=13, ha="center", va="center", fontweight="bold", color="blue")
    ax.text(13, 8, f"Scale X = {resize_result.scale_x:.4f}\n"
                    f"Scale Y = {resize_result.scale_y:.4f}\n"
                    f"Pixels/output = {resize_result.pixels_per_output:.2f}",
            fontsize=10, ha="center", family="monospace",
            bbox=dict(boxstyle="round", facecolor="lightyellow"))
    images["s4_concept"] = _save(fig, "s4_concept.png")

    # 4-c  zoom comparison
    fig, axes = plt.subplots(1, 2, figsize=(14, 6))
    cy, cx = warped.shape[0]//2, warped.shape[1]//2
    axes[0].imshow(_to8(warped[cy-50:cy+50, cx-50:cx+50]), cmap="gray", interpolation="nearest")
    axes[0].set_title("Warped Center 100×100\n(camera pixels)", fontweight="bold")
    cy2, cx2 = resized.shape[0]//2, resized.shape[1]//2
    axes[1].imshow(_to8(resized[cy2-15:cy2+15, cx2-15:cx2+15]), cmap="gray", interpolation="nearest")
    axes[1].set_title("Resized Center 30×30\n(display pixels)", fontweight="bold")
    plt.tight_layout()
    images["s4_zoom"] = _save(fig, "s4_zoom.png")

    # ── Step 5  16-bit normalisation (6_image_normalizer.py) ──────
    print("\n[Step 5] 16-bit normalisation …")
    normalizer = ImageNormalizer(bit_depth=16)
    norm_result = normalizer.normalize(resized)
    normalized = norm_result.image
    print(f"  input range  : [{norm_result.original_min:.2f}, {norm_result.original_max:.2f}]")
    print(f"  output range : [{norm_result.normalized_min}, {norm_result.normalized_max}]")

    # 5-a  histograms + result
    fig, axes = plt.subplots(2, 2, figsize=(16, 12))
    axes[0, 0].hist(resized.ravel(), bins=256, color="steelblue", alpha=0.8)
    axes[0, 0].set_title("Before Normalisation", fontweight="bold")
    axes[0, 0].set_xlabel("Pixel Value"); axes[0, 0].set_ylabel("Frequency")
    axes[0, 1].hist(normalized.ravel(), bins=256, color="coral", alpha=0.8)
    axes[0, 1].set_title("After 16-bit Normalisation (0–65 535)", fontweight="bold")
    axes[0, 1].set_xlabel("Pixel Value (16-bit)")
    axes[1, 0].imshow(_to8(normalized), cmap="gray")
    axes[1, 0].set_title(f"Normalised {normalized.shape[1]}×{normalized.shape[0]}  |  "
                         f"{normalized.dtype}", fontweight="bold")
    ax = axes[1, 1]; ax.axis("off")
    info = (f"Input  Range : [{norm_result.original_min:.4f}, {norm_result.original_max:.4f}]\n"
            f"Output Range : [{norm_result.normalized_min}, {norm_result.normalized_max}]\n"
            f"Bit Depth    : {norm_result.bit_depth}-bit  (max = {2**norm_result.bit_depth - 1})\n\n"
            f"Formula:\n  out = (px − min) / (max − min) × 65 535")
    ax.text(0.1, 0.7, info, fontsize=13, transform=ax.transAxes, va="top",
            family="monospace",
            bbox=dict(boxstyle="round,pad=0.5", facecolor="#FFF3E0"))
    plt.tight_layout()
    images["s5_norm"] = _save(fig, "s5_norm.png")

    # ── Step 6  Pipeline summary ──────────────────────────────────
    print("\n[Step 6] Summary visualisation …")

    # 6-a  5-stage strip
    fig, axes = plt.subplots(1, 5, figsize=(26, 4.5))
    stages = [
        (original,   "Original"),
        (cropped,    "Cropped"),
        (warped,     "Warped"),
        (resized,    "Resized"),
        (normalized, "Normalised"),
    ]
    for i, (img, title) in enumerate(stages):
        axes[i].imshow(_to8(img), cmap="gray")
        axes[i].set_title(f"{title}\n{img.shape[1]}×{img.shape[0]}",
                          fontsize=10, fontweight="bold")
        axes[i].set_xticks([]); axes[i].set_yticks([])
    plt.tight_layout()
    images["s6_summary"] = _save(fig, "s6_summary.png")

    # 6-b  architecture diagram
    fig, ax = plt.subplots(figsize=(14, 7))
    ax.axis("off"); ax.set_xlim(0, 14); ax.set_ylim(0, 9)
    ax.set_title("Pipeline Architecture (src/ modules)", fontsize=16, fontweight="bold", pad=20)
    modules = [
        (2, 7, "1_utils.py",           "Utility\nFunctions",    "#E3F2FD"),
        (2, 5, "3_image_cropper.py",   "Image\nCrop",          "#FFF3E0"),
        (5.5, 5, "2_roi_detector.py",  "ROI\nDetection",       "#E8F5E9"),
        (9, 7, "4_perspective_warper.py","Perspective\nWarp",   "#FCE4EC"),
        (9, 5, "5_area_sum_resizer.py", "Area-Sum\nResize",    "#E0F7FA"),
        (9, 3, "6_image_normalizer.py", "16-bit\nNormalise",   "#F3E5F5"),
        (5.5, 3, "7_display_panel_\nprocessor.py","Pipeline\nOrchestrator","#E8EAF6"),
    ]
    for x, y, name, desc, col in modules:
        ax.add_patch(plt.Rectangle((x-1.3, y-0.65), 2.6, 1.3,
                     lw=1.5, ec="#333", fc=col, alpha=0.9, zorder=2))
        ax.text(x, y+0.2, name, fontsize=7.5, ha="center", va="center",
                fontweight="bold", zorder=3)
        ax.text(x, y-0.25, desc, fontsize=7, ha="center", va="center",
                color="#555", zorder=3)
    arrp = dict(arrowstyle="->", color="#1877F2", lw=2)
    ax.annotate("", xy=(2, 5.65),   xytext=(2, 6.35),   arrowprops=arrp)
    ax.annotate("", xy=(4.2, 5),    xytext=(3.2, 5),    arrowprops=arrp)
    ax.annotate("", xy=(7.7, 7),    xytext=(6.8, 5.5),  arrowprops=arrp)
    ax.annotate("", xy=(9, 6.35),   xytext=(9, 5.65),   arrowprops=arrp)
    ax.annotate("", xy=(9, 3.65),   xytext=(9, 4.35),   arrowprops=arrp)
    images["s6_arch"] = _save(fig, "s6_arch.png")

    print(f"\nDone: All {len(images)} images saved to {IMG_DIR}")
    return images, config, roi, warp_result, resize_result, norm_result


# ═══════════════════════════════════════════════════════════════════
# PPT builder helpers
# ═══════════════════════════════════════════════════════════════════
def _add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(slide, left, top, width, height, text,
              size=18, bold=False, color=GRAY, align=PP_ALIGN.LEFT, font_name="Calibri"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(size); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = font_name; p.alignment = align
    return box


def _add_bullets(slide, left, top, width, height, lines,
                 size=13, color=GRAY):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = box.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(size); p.font.color.rgb = color
        p.font.name = "Calibri"; p.space_after = Pt(4)
    return box


def _add_img(slide, path, left, top, width=None, height=None):
    kw = {}
    if width:  kw["width"]  = Inches(width)
    if height: kw["height"] = Inches(height)
    if Path(path).exists():
        return slide.shapes.add_picture(str(path), Inches(left), Inches(top), **kw)
    return None


def _section_header(prs, slide_layouts, text):
    """Full-width blue section header slide."""
    slide = prs.slides.add_slide(slide_layouts)
    _add_bg(slide, BLUE)
    _add_text(slide, 1, 2.8, 11, 1.5, text,
              size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    return slide


# ═══════════════════════════════════════════════════════════════════
# Build PPT
# ═══════════════════════════════════════════════════════════════════
def build_ppt(images, config, roi, warp_result, resize_result, norm_result):
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── Slide 1  Title ────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _add_bg(s, BLUE)
    _add_text(s, 1, 1.5, 11, 1.5,
              "Display Panel ROI Algorithm",
              size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_text(s, 1, 3.2, 11, 1,
              "Step-by-Step Pipeline Results with Intermediate Images",
              size=24, color=LTBLUE, align=PP_ALIGN.CENTER)
    _add_text(s, 1, 5.0, 11, 0.8,
              "Code Order: 1_utils → 3_image_cropper → 2_roi_detector → "
              "4_perspective_warper → 5_area_sum_resizer → 6_image_normalizer",
              size=14, color=MIDBLUE, align=PP_ALIGN.CENTER)
    _add_text(s, 1, 6.2, 11, 0.5,
              "MR Display Hardware Team  |  Byung Geun (BG) Jun  |  March 2026",
              size=14, color=DKBLUE, align=PP_ALIGN.CENTER)

    # ── Slide 2  Pipeline Overview ────────────────────────────────
    s = prs.slides.add_slide(blank)
    _add_text(s, 0.5, 0.3, 12, 0.8,
              "Pipeline Overview", size=32, bold=True, color=BLUE)
    _add_bullets(s, 0.5, 1.2, 5.5, 5.5, [
        "End-to-End Pipeline Flow:",
        "",
        "Step 0 — Load Original Image (1_utils.py)",
        "  • Read 16-bit TIFF camera image",
        "",
        "Step 1 — Image Cropping (3_image_cropper.py)",
        "  • Remove unnecessary border regions",
        "",
        "Step 2 — ROI Detection (2_roi_detector.py)",
        "  • Binarisation → Morphology → Contour → Corners",
        "",
        "Step 3 — Perspective Warp (4_perspective_warper.py)",
        "  • 4-corner perspective transform",
        "",
        "Step 4 — Area-Sum Resize (5_area_sum_resizer.py)",
        "  • Camera resolution → Display resolution",
        "",
        "Step 5 — 16-bit Normalisation (6_image_normalizer.py)",
        "  • Scale to full dynamic range (0–65 535)",
    ], size=12)
    _add_img(s, images["s6_arch"], 6.5, 1.2, width=6.3)

    # ── Slide 3  Step 0: Original Image ──────────────────────────
    _section_header(prs, blank, "Step 0 — Load Original Image\n(1_utils.py)")

    s = prs.slides.add_slide(blank)
    _add_text(s, 0.5, 0.3, 12, 0.8,
              "Step 0: Original Image", size=32, bold=True, color=BLUE)
    _add_img(s, images["s0_original"], 0.3, 1.1, width=12.5)

    # ── Slide 4  Step 0: Histogram ───────────────────────────────
    s = prs.slides.add_slide(blank)
    _add_text(s, 0.5, 0.3, 12, 0.8,
              "Step 0: Image Statistics & Histogram", size=32, bold=True, color=BLUE)
    _add_img(s, images["s0_histogram"], 0.3, 1.1, width=12.5)
    _add_bullets(s, 0.5, 5.8, 12, 1.5, [
        "• Histogram shows background (low values) vs. display area (high values)",
        "• Red dashed line: binarisation threshold used for ROI detection",
    ], size=12)

    # ── Slide 5  Step 1: Image Cropping ──────────────────────────
    _section_header(prs, blank, "Step 1 — Image Cropping\n(3_image_cropper.py)")

    s = prs.slides.add_slide(blank)
    _add_text(s, 0.5, 0.3, 12, 0.8,
              "Step 1: Crop — Before & After", size=32, bold=True, color=BLUE)
    _add_img(s, images["s1_crop"], 0.3, 1.1, width=12.5)
    _add_bullets(s, 0.5, 5.6, 12, 1.8, [
        "■ Classes: ImageCropper + CropRegion",
        f"  • Crop region: ({config.crop_x_start}, {config.crop_y_start}) → "
        f"({config.crop_x_end}, {config.crop_y_end})",
        "  • Implementation: numpy slicing — image[y_start:y_end, x_start:x_end]",
        "  • Purpose: Speed up ROI detection + eliminate border noise",
    ], size=12)

    # ── Slide 6  Step 2: ROI Detection — Binarisation ────────────
    _section_header(prs, blank, "Step 2 — ROI Detection\n(2_roi_detector.py)")

    s = prs.slides.add_slide(blank)
    _add_text(s, 0.5, 0.3, 12, 0.8,
              "Step 2a: Binarisation & Morphology", size=28, bold=True, color=BLUE)
    _add_img(s, images["s2_binary"], 0.3, 1.1, width=12.5)
    _add_bullets(s, 0.5, 5.3, 12, 2, [
        "■ Class: ROIDetector",
        f"  1) Binary threshold: pixel > {config.roi_threshold} → white",
        f"  2) Morphological Close ({config.morph_kernel_size}×{config.morph_kernel_size} kernel): fill holes",
        f"  3) Morphological Open ({config.morph_kernel_size}×{config.morph_kernel_size} kernel): remove noise",
        "  4) scipy.ndimage.binary_fill_holes(): fill remaining voids",
    ], size=12)

    # ── Slide 7  Step 2: Contour & Corners ───────────────────────
    s = prs.slides.add_slide(blank)
    _add_text(s, 0.5, 0.3, 12, 0.8,
              "Step 2b: Contour Detection & Corner Extraction", size=28, bold=True, color=BLUE)
    _add_img(s, images["s2_corners"], 0.3, 1.1, width=8.5)
    _add_img(s, images["s2_table"], 9.0, 1.1, width=4)
    _add_bullets(s, 0.5, 5.3, 12, 2, [
        "■ Algorithm:",
        "  1) cv2.findContours() → extract external contours from binary mask",
        "  2) max(contours, key=cv2.contourArea) → select largest contour",
        "  3) cv2.minAreaRect() + cv2.boxPoints() → compute 4 corners",
        "  4) Corner ordering: sort by sum/diff → TL, TR, BR, BL",
        f"  5) Tilt angle: {roi.angle:.4f}° — arctan2(top_edge_dy, top_edge_dx)",
    ], size=12)

    # ── Slide 8  Step 3: Perspective Warp ────────────────────────
    _section_header(prs, blank, "Step 3 — Perspective Correction\n(4_perspective_warper.py)")

    s = prs.slides.add_slide(blank)
    _add_text(s, 0.5, 0.3, 12, 0.8,
              "Step 3a: Warp — Before & After", size=32, bold=True, color=BLUE)
    _add_img(s, images["s3_warp"], 0.3, 1.1, width=8)
    _add_bullets(s, 8.5, 1.2, 4.5, 4, [
        "■ Class: PerspectiveWarper",
        "",
        "Core Operations:",
        "  1) Map 4 source corners → 4 dest corners",
        "  2) cv2.getPerspectiveTransform()",
        "     → Compute 3×3 homography matrix",
        "  3) cv2.warpPerspective()",
        "     → Apply the perspective transform",
        "",
        "■ Interpolation: INTER_NEAREST",
        "  • Nearest-neighbor interpolation",
        "  • Preserves original pixel values",
        "  • Critical for measurement accuracy",
    ], size=11)

    # ── Slide 9  Step 3: Transform Matrix ────────────────────────
    s = prs.slides.add_slide(blank)
    _add_text(s, 0.5, 0.3, 12, 0.8,
              "Step 3b: Transform Matrix & Corner Mapping", size=28, bold=True, color=BLUE)
    _add_img(s, images["s3_matrix"], 1.5, 1.1, width=10)
    _add_bullets(s, 0.5, 5.5, 12, 1.5, [
        "• 3×3 homography matrix maps tilted quadrilateral to axis-aligned rectangle",
        "• INTER_NEAREST ensures original measurement values are exactly preserved",
    ], size=12)

    # ── Slide 10  Step 4: Area-Sum Resize ────────────────────────
    _section_header(prs, blank, "Step 4 — Area-Sum Resize\n(5_area_sum_resizer.py)")

    s = prs.slides.add_slide(blank)
    _add_text(s, 0.5, 0.3, 12, 0.8,
              "Step 4a: Camera → Display Resolution", size=28, bold=True, color=BLUE)
    _add_img(s, images["s4_resize"], 0.3, 1.1, width=12.5)
    _add_bullets(s, 0.5, 5.3, 12, 2, [
        "■ Class: AreaSumResizer — camera resolution → display resolution conversion",
        f"  • Scale: {resize_result.scale_x:.4f} × {resize_result.scale_y:.4f} "
        f"(camera px / display px)",
        f"  • Output: {config.display_width} × {config.display_height} "
        "(1:1 mapping to actual display)",
        "  • Method: weighted area-based summation of all contributing camera pixels",
    ], size=12)

    # ── Slide 11  Step 4: Concept ────────────────────────────────
    s = prs.slides.add_slide(blank)
    _add_text(s, 0.5, 0.3, 12, 0.8,
              "Step 4b: Area-Sum Concept Diagram", size=28, bold=True, color=BLUE)
    _add_img(s, images["s4_concept"], 0.3, 1.1, width=7)
    _add_bullets(s, 7.5, 1.2, 5.3, 5, [
        "■ Area-Sum Resize Principle:",
        "",
        "In the camera image, one display pixel",
        f"corresponds to ~{resize_result.scale_x:.1f}×{resize_result.scale_y:.1f}",
        "camera pixels.",
        "",
        "The output value for each display pixel",
        "is the mean of all camera pixels within",
        "the corresponding area.",
        "",
        "■ Why Area-Sum?",
        "  • Standard resize (bilinear, etc.) alters",
        "    original values via interpolation",
        "  • Area-Sum directly averages all pixels",
        "    → preserves measurement accuracy",
        "  • Essential for precise sub-pixel",
        "    brightness in display measurement",
    ], size=11)

    # ── Slide 12  Step 4: Zoom ───────────────────────────────────
    s = prs.slides.add_slide(blank)
    _add_text(s, 0.5, 0.3, 12, 0.8,
              "Step 4c: Zoom — Camera vs. Display Pixels", size=28, bold=True, color=BLUE)
    _add_img(s, images["s4_zoom"], 0.3, 1.1, width=12.5)
    _add_bullets(s, 0.5, 5.5, 12, 1.5, [
        "• Left: Warped image center 100×100 — individual camera pixels visible",
        "• Right: Resized image center 30×30 — each pixel = one display pixel",
    ], size=12)

    # ── Slide 13  Step 5: Normalisation ──────────────────────────
    _section_header(prs, blank, "Step 5 — 16-bit Normalisation\n(6_image_normalizer.py)")

    s = prs.slides.add_slide(blank)
    _add_text(s, 0.5, 0.3, 12, 0.8,
              "Step 5: Normalisation — Histograms & Result", size=28, bold=True, color=BLUE)
    _add_img(s, images["s5_norm"], 0.3, 1.1, width=12.5)

    # ── Slide 14  Full Pipeline Summary ──────────────────────────
    _section_header(prs, blank, "Pipeline Summary")

    s = prs.slides.add_slide(blank)
    _add_text(s, 0.5, 0.3, 12, 0.8,
              "Full Pipeline — All 5 Stages", size=32, bold=True, color=BLUE)
    _add_img(s, images["s6_summary"], 0.3, 1.2, width=12.5)
    _add_bullets(s, 0.5, 4.8, 12, 2.5, [
        "Original → Crop → ROI Detect → Perspective Warp → Area-Sum Resize → 16-bit Normalise",
        "",
        "■ Key Design Principles:",
        "  • INTER_NEAREST interpolation: preserves original pixel values",
        "  • Area-Sum resize: every camera pixel is accounted for",
        "  • Modular OOP architecture: each step is an independent class",
        "  • Orchestrated by 7_display_panel_processor.py (DisplayPanelProcessor)",
    ], size=13)

    # ── Slide 15  Architecture ───────────────────────────────────
    s = prs.slides.add_slide(blank)
    _add_text(s, 0.5, 0.3, 12, 0.8,
              "Code Architecture (src/ Directory)", size=32, bold=True, color=BLUE)
    _add_img(s, images["s6_arch"], 0.3, 1.2, width=6.5)
    _add_bullets(s, 7, 1.2, 6, 5.5, [
        "src/",
        "├── __init__.py              Package init",
        "├── 1_utils.py               Utilities",
        "├── 2_roi_detector.py        ROI detection",
        "│     ├── ROIDetector",
        "│     └── AdaptiveROIDetector",
        "├── 3_image_cropper.py       Image cropping",
        "│     ├── CropRegion",
        "│     └── ImageCropper",
        "├── 4_perspective_warper.py   Perspective warp",
        "│     └── PerspectiveWarper",
        "├── 5_area_sum_resizer.py    Area-sum resize",
        "│     └── AreaSumResizer",
        "├── 6_image_normalizer.py    Normalisation",
        "│     └── ImageNormalizer",
        "└── 7_display_panel_processor.py",
        "      ├── ProcessingConfig",
        "      ├── ProcessingResult",
        "      └── DisplayPanelProcessor",
    ], size=10)

    # ── Slide 16  Thank You ──────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _add_bg(s, BLUE)
    _add_text(s, 1, 2.5, 11, 1.5,
              "Thank You", size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_text(s, 1, 4.2, 11, 1,
              "ROI Algorithm — Step-by-Step Pipeline Documentation",
              size=20, color=LTBLUE, align=PP_ALIGN.CENTER)
    _add_text(s, 1, 5.5, 11, 0.5,
              "Generated from src/ modules  |  MR Display Hardware Team",
              size=14, color=MIDBLUE, align=PP_ALIGN.CENTER)

    # ── Save ─────────────────────────────────────────────────────
    ppt_path = DOCS_DIR / "ROI_Pipeline_StepByStep.pptx"
    prs.save(str(ppt_path))
    print(f"\nPPT saved: {ppt_path}")
    print(f"  Total slides: {len(prs.slides)}")
    return str(ppt_path)


# ═══════════════════════════════════════════════════════════════════
# Main
# ═══════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    imgs, cfg, roi, wr, rr, nr = run_pipeline()
    build_ppt(imgs, cfg, roi, wr, rr, nr)
